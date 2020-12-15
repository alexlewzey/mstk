"""set of tools for creating and manipulating pptx files."""

from typing import Tuple, Optional, List, Dict

import numpy as np
import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.slide import Slide
from pptx.table import Table
from pptx.util import Cm, Pt

SLIDE_WIDTH: float = 26
SLIDE_HEIGHT: float = 15
TOP = 1.6
LEFT = 0.1
HEIGHT = SLIDE_HEIGHT - TOP
WIDTH = SLIDE_WIDTH - (2 * LEFT)

HEADLINE_HEIGHT: float = 0.1
HEADLINE_TOP: float = 1.2
HEADLINE_LEFT: float = 0

LOGO_HEIGHT: float = 1
LOGO_TOP: float = 0.1
LOGO_LEFT: float = SLIDE_WIDTH - 1.1

TITLE_HEIGHT: float = 1
TITLE_WIDTH: float = SLIDE_WIDTH - 2
TITLE_TOP: float = 0.1
TITLE_LEFT: float = 0.3

colors = {
    'green_main': (112, 182, 88),
    'green_dark': (33, 84, 37),
    'grey_dark': (49, 45, 49),
    'dog': (191, 209, 67),
    'cat': (232, 132, 65),
    'small_pet': (212, 153, 59),
    'fish': (40, 58, 140),
    'bird': (109, 173, 218),
    'reptile': (101, 38, 57),
    'black': (0, 0, 0),
    'white': (255, 255, 255),
}


def add_slide_text_and_image(prs: Presentation, text: str, title: str, img_path: str, path_logo: str,
                             proportions: Tuple[float, ...]) -> Tuple[Presentation, Slide]:
    prs, slide, dims = add_slide_text_and_blank(prs=prs, text=text, title=title, path_logo=path_logo,
                                                proportions=proportions)
    slide = add_image(slide=slide, img_path=img_path, left=dims['left1'], top=TOP, width=dims['width1'])
    return prs, slide


def add_slide_text_and_blank(prs: Presentation, text: str, title: str, path_logo: str,
                             proportions: Tuple[float, ...]) -> Tuple[Presentation, Slide, Dict]:
    dims = make_col_dims(left=LEFT, width=WIDTH, proportions=proportions)
    prs, slide = add_new_slide(prs, title=title, path_logo=path_logo)
    slide = add_textframe(slide=slide, left=dims['left0'], top=TOP, width=dims['width0'], height=HEIGHT, text=text)
    return prs, slide, dims


def make_col_dims(left: float, width: float, proportions: Tuple[float, ...], top: float = TOP,
                  height: float = HEIGHT) -> Dict[
    str, float]:
    """split a rectangle in separate columns a specified proportion and return the dimensions of the subsequent columns
    """
    dims: Dict = {}
    ncols = len(proportions)
    lefts, widths = _make_sub_dims(start=left, length=width, proportions=proportions)

    for i in range(ncols):
        dims[f'left{i}'] = lefts[i]
        dims[f'width{i}'] = widths[i]
        dims[f'top{i}'] = top
        dims[f'height{i}'] = height
    return dims


def make_row_dims(top: float, height: float, proportions: Tuple[float, ...], left: float = LEFT,
                  width: float = WIDTH) -> Dict[str, float]:
    """split a rectangle in separate rows a specified proportion and return the dimensions of the subsequent rows"""
    dims: Dict = {}
    nrows = len(proportions)
    tops, heights = _make_sub_dims(start=top, length=height, proportions=proportions)

    for i in range(nrows):
        dims[f'left{i}'] = left
        dims[f'width{i}'] = width
        dims[f'top{i}'] = tops[i]
        dims[f'height{i}'] = heights[i]
    return dims


def _make_sub_dims(start, length, proportions):
    if sum(proportions) != 1.0:
        raise ValueError(f'Proportions should sum to 1: {sum(proportions)}')

    proportions = np.array(proportions)
    proportions_cum = np.cumsum(proportions)
    starts = [start] + list(((proportions_cum * length) + start)[:-1])
    lengths = proportions * length
    return starts, lengths


def add_new_slide(prs: Presentation, title: str, path_logo: str) -> Tuple[Presentation, Slide]:
    """add a new slide with logo and title returning the presentation and the new slide"""
    prs, slide = add_blank_slide(prs)
    slide = add_title(slide, title)
    slide = add_logo(slide, path_logo)
    slide = add_header_line(slide)
    return prs, slide


def add_title(slide: Slide, title: str) -> Slide:
    """add formatted title to slide"""
    return add_textframe(slide=slide, left=TITLE_LEFT, top=TITLE_TOP, width=TITLE_WIDTH, height=TITLE_HEIGHT,
                         text=title,
                         rgb=colors['green_main'])


def add_logo(slide, path_logo: str) -> Slide:
    """add logo to slide"""
    slide = add_image(slide=slide, img_path=path_logo, left=LOGO_LEFT, top=LOGO_TOP, height=LOGO_HEIGHT)
    return slide


def add_header_line(slide) -> Slide:
    """add title underline (big green strip across the slide)"""
    return add_shape(slide=slide, left=HEADLINE_LEFT, top=HEADLINE_TOP, width=SLIDE_WIDTH, height=HEADLINE_HEIGHT,
                     rgb=colors['green_main'])


def add_blank_slide(prs: Presentation) -> Tuple[Presentation, Slide]:
    """add blank slide to the end of the presentation"""
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    return prs, slide


def add_textframe(slide: Slide, left: float, top: float, width: float, height: float, text: str,
                  rgb: Tuple[int, int, int] = (0, 0, 0)) -> Slide:
    """
    add a text box in specified location
    Args:
        slide: slide object
        left: horizontal position cm
        top: vertical position cm
        width: cm
        height: cm
        text: text content of frame
        rgb: color of text in rgb

    Returns:
        slide and textframe object
    """
    left, top, width, height = _get_cm(left, top, width, height)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    font = run.font
    font.color.rgb = RGBColor(*rgb)
    return slide


def add_image(slide: Slide, img_path: str, left: float, top: float, width=None, height=None) -> Slide:
    """
    add image to slide in a specified location
    Args:
        slide: slide object
        img_path: path to image file png or jpg
        left: horizontal position cm
        top: vertical position cm
        width: cm
        height: cm

    Returns:
        slide and picture objects
    """
    left, top, width, height = _get_cm(left, top, width, height)
    slide.shapes.add_picture(img_path, left, top, height=height, width=width)
    return slide


def add_shape(slide: Slide, left: float, top: float, width: float, height: float, rgb: Tuple[int, int, int],
              shape_type=MSO_SHAPE.RECTANGLE) -> Slide:
    """
    add a filled shape to a specified location on the slide
    Args:
        slide: pptx slide object
        left: horizontal position cm
        top: vertical position cm
        width: cm
        height: cm
        rgb: color of shape in rgb
        shape_type: pptx shape type object

    Returns:
        the slide
    """
    left, top, width, height = _get_cm(left, top, width, height)
    shape = slide.shapes
    box = shape.add_shape(shape_type, left, top, width, height)
    fill = box.fill
    line = box.line
    line.color.rgb = RGBColor(*rgb)
    fill.solid()
    fill.fore_color.rgb = RGBColor(*rgb)
    return slide


def add_table_from_df(slide: Slide, df: pd.DataFrame, left: float, top: float, width: float, height: float,
                      col_proportions: Optional[List[float]] = None, first_cell: str = ' ',
                      font_size: float = 12) -> Slide:
    """
    transform a df into a pptx table and add it to the passed slide in a specified location
    Args:
        slide: pptx slide object
        df: DataFrame table
        left: horizontal position cm
        top: vertical position cm
        width: cm
        height: cm
        col_proportions: array like with width proportion of each column

    Returns:
        slide and table objects
    """
    rows = df.index.tolist()
    columns = df.columns.tolist()
    values = [list(row) for row in df.values]
    return _add_table_from_lst(slide=slide, cols=columns, rows=rows, values=values, left=left, top=top, width=width,
                               height=height, col_proportions=col_proportions, first_cell=first_cell, font_size=12)


def _get_cm(left: float, top: float, width: Optional[float], height: Optional[float]) -> Tuple[
    Cm, Cm, Optional[Cm], Optional[Cm]]:
    """convert floats to pptx centimeter objects"""
    left = Cm(left)
    top = Cm(top)
    width = Cm(width) if width else None
    height = Cm(height) if height else None
    return left, top, width, height


def set_slide_size(prs: Presentation, width: float = 26, height: float = 15) -> Presentation:
    """set the size of all slides in cm"""
    prs.slide_width = Cm(width)
    prs.slide_height = Cm(height)
    return prs


def _add_table_from_lst(slide: Slide, cols: List[str], rows: List[str], values: List[List[str]], left: float,
                        top: float,
                        width: float, height: float, first_cell: str = ' ', font_size: float = 12,
                        col_proportions: Optional[List[float]] = None) -> Slide:
    """convert python list objects to pptx table and add it to the passed slide in a specified location"""
    nrows, ncols = len(rows) + 1, len(cols) + 1
    left_cm, top_cm, width_cm, height_cm = _get_cm(left, top, width, height)
    shapes = slide.shapes
    table = shapes.add_table(nrows, ncols, left_cm, top_cm, width_cm, height_cm).table

    # set column widths
    assert sum(col_proportions) == 1.
    if col_proportions:
        for i, proportion in enumerate(col_proportions):
            table.columns[i].width = Cm(width * proportion)

    table = _populate_table(table=table, cols=cols, rows=rows, values=values, first_cell=first_cell)

    if font_size:
        _set_table_font_size(table, font_size)

    return slide


def _populate_table(table: Table, cols: List[str], rows: List[str], values: List[List[str]],
                    first_cell: str = ' ') -> Table:
    """populate the cells of a pptx table from lists"""
    table.cell(0, 0).text = str(first_cell)
    for j, col in enumerate(cols):
        table.cell(0, j + 1).text = str(col)
    for i, row in enumerate(rows):
        table.cell(i + 1, 0).text = str(row)
    for j in range(len(cols)):
        for i in range(len(rows)):
            table.cell(i + 1, j + 1).text = str(values[i][j])
    return table


def _set_table_font_size(table: Table, pt: float = 12) -> Table:
    """set the font size of every cell in a table object"""
    for cell in _iter_cells(table):
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(pt)
    return table


def _iter_cells(table: Table):
    """generator for every cell in a table"""
    for row in table.rows:
        for cell in row.cells:
            yield cell
